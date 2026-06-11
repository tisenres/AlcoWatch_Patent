#!/usr/bin/env python3
"""Build Major Project defense presentation by editing the NTCC template.

Preserves design 1:1 (fonts, colors, layouts, decorative shapes) — replaces
text content slide-by-slide with Major Project (driver stress detection)
material. Swaps the BAC architecture image on slide 8 and the BAC hardware
image on slide 10 with the stress paper's system_overview.png and
model_architecture.png.

Run:
    python3 scripts/build_major_project_pptx.py
Output:
    Anastasiia_Major_Project_Defense.pptx (in repo root)
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

GITHUB_URL = "https://github.com/tisenres/AlcoWatch_Patent"

REPO_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = REPO_ROOT / "Anastasiia_NTCC (2).pptx"
OUTPUT = REPO_ROOT / "Anastasiia_Major_Project_Defense.pptx"
SYSTEM_OVERVIEW = REPO_ROOT / "paper_figures" / "stress" / "system_overview.png"
MODEL_ARCH = REPO_ROOT / "paper_figures" / "stress" / "model_architecture.png"


def replace_paragraph_text(paragraph, new_text):
    if not paragraph.runs:
        paragraph.add_run().text = new_text
        return
    paragraph.runs[0].text = new_text
    for extra in paragraph.runs[1:]:
        extra._r.getparent().remove(extra._r)


def set_text_frame_lines(text_frame, lines):
    if not lines or not text_frame.paragraphs:
        return
    template_p = text_frame.paragraphs[0]
    template_run = template_p.runs[0] if template_p.runs else None

    template_rpr = None
    if template_run is not None:
        rpr = template_run._r.find(qn("a:rPr"))
        if rpr is not None:
            template_rpr = copy.deepcopy(rpr)

    template_ppr = template_p._p.find(qn("a:pPr"))
    if template_ppr is not None:
        template_ppr = copy.deepcopy(template_ppr)

    for p in list(text_frame.paragraphs[1:]):
        p._p.getparent().remove(p._p)

    replace_paragraph_text(template_p, lines[0])

    for line in lines[1:]:
        new_p = text_frame.add_paragraph()
        if template_ppr is not None:
            new_p._p.insert(0, copy.deepcopy(template_ppr))
        run = new_p.add_run()
        run.text = line
        if template_rpr is not None:
            run._r.insert(0, copy.deepcopy(template_rpr))


def find_shape(slide, name):
    return next((s for s in slide.shapes if s.name == name), None)


def set_shape_text(slide, name, lines, autofit=False):
    shape = find_shape(slide, name)
    if shape is None:
        print(f"  WARN: shape {name!r} not found", file=sys.stderr)
        return
    if not shape.has_text_frame:
        print(f"  WARN: shape {name!r} has no text frame", file=sys.stderr)
        return
    if isinstance(lines, str):
        lines = [lines]
    set_text_frame_lines(shape.text_frame, lines)
    if autofit:
        # Tells PowerPoint to shrink text proportionally on overflow.
        # Safety net for long bullets that don't fit at native font size.
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        shape.text_frame.word_wrap = True


def replace_in_runs(slide, name, replacements):
    """In-place find/replace per run — preserves tabs, blank paragraphs, fonts."""
    shape = find_shape(slide, name)
    if shape is None or not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            for old, new in replacements:
                text = text.replace(old, new)
            if text != run.text:
                run.text = text


def add_centered_textbox(slide, left, top, width, height, lines, font_size=18,
                         bold=False, color=(60, 60, 60), font_name="Arial"):
    """Add a fresh textbox with centered text. Used for GitHub link callout."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
    return tb


def replace_picture(slide, name, new_image_path):
    target = find_shape(slide, name)
    if target is None:
        print(f"  WARN: picture {name!r} not found", file=sys.stderr)
        return
    if target.shape_type != MSO_SHAPE_TYPE.PICTURE:
        print(f"  WARN: shape {name!r} is not a picture", file=sys.stderr)
        return
    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)
    slide.shapes.add_picture(str(new_image_path), left, top, width=width, height=height)


def main():
    if not TEMPLATE.exists():
        sys.exit(f"ERROR: template not found: {TEMPLATE}")
    if not SYSTEM_OVERVIEW.exists() or not MODEL_ARCH.exists():
        sys.exit("ERROR: stress figures not found")

    print(f"Loading: {TEMPLATE.name}")
    p = Presentation(str(TEMPLATE))
    slides = list(p.slides)
    if len(slides) != 15:
        sys.exit(f"ERROR: expected 15 slides, got {len(slides)}")

    # ---- Slide 1: Title ----
    print("Slide  1: Title")
    s = slides[0]
    set_shape_text(s, "Title 1",
        "AI-BASED REAL-TIME DRIVER STRESS DETECTION AND "
        "ADAPTIVE VEHICLE ENVIRONMENT CONTROL SYSTEM"
    )
    replace_in_runs(s, "TextBox 4", [
        ("Semester 5", "Semester 6"),
        ("ETMN100", "ETMJ100"),
        ("Minor Project", "Major Project"),
    ])
    # TextBox 5 split as runs ['1', 'st'(superscript), '  NTCC Presentation '] —
    # whole-text rewrite drops the superscript and inherits regular run formatting.
    set_shape_text(s, "TextBox 5", "Final Presentation")

    # ---- Slide 2: Table of Contents ----
    print("Slide  2: Table of Contents")
    s = slides[1]
    set_shape_text(s, "Title 1", "Table of Content")
    set_shape_text(s, "Content Placeholder 4", [
        "Synopsis",
        "Introduction",
        "Literature Review",
        "Limitations of Existing Systems",
        "Proposed Solution",
        "System Architecture",
        "Methodology & Model",
        "Results & Evaluation",
        "Conclusion & References",
    ])

    # ---- Slide 3: Synopsis ----
    print("Slide  3: Synopsis")
    s = slides[2]
    set_shape_text(s, "Title 1", "Synopsis")
    set_shape_text(s, "Text 2", "Objective:")
    set_shape_text(s, "Text 3",
        "AI-based real-time driver stress detection on a wrist-only wearable, "
        "paired with adaptive vehicle cabin control to mitigate stress-induced "
        "crash risk via graduated, non-punitive interventions.",
        autofit=True,
    )
    set_shape_text(s, "Text 4", "Key Innovation:")
    set_shape_text(s, "Text 5",
        "Closed-loop physiological sensing → on-device BiLSTM + attention "
        "(4-class) → 12-byte BLE → graduated cabin response (LED, fan, audio) "
        "on Arduino Nano 33 BLE.",
        autofit=True,
    )

    # ---- Slide 4: Introduction ----
    print("Slide  4: Introduction")
    s = slides[3]
    set_shape_text(s, "Title 1", "Introduction")
    set_shape_text(s, "Text 2", "Current Challenge:")
    set_shape_text(s, "Text 3", [
        "Road crashes kill 1.19M people annually (WHO); ~90% from human error",
        "Acute stress raises crash risk 2–3× via HRV drop, slow reactions, impulsive maneuvers",
        "Modern vehicles handle perceptual lapses but ignore autonomic state",
    ], autofit=True)
    set_shape_text(s, "Text 4", "Need for Innovation:")
    set_shape_text(s, "Text 5",
        "Consumer wearables now ship sensor-grade PPG/EDA/TEMP/ACC; on-device "
        "TFLite enables real-time inference. Closed-loop cabin adaptation can "
        "reduce stress without driver intervention.",
        autofit=True,
    )

    # ---- Slide 5: Literature Review ----
    print("Slide  5: Literature Review")
    s = slides[4]
    set_shape_text(s, "Title 1", "Literature Review")
    set_shape_text(s, "Text 2", "Physiological Stress Markers:")
    set_shape_text(s, "Text 3", [
        "HRV (PPG): RMSSD/SDNN drop; LF/HF balance shifts under stress",
        "EDA: tonic level rises, phasic peaks denser under sympathetic activation",
        "Skin temperature: peripheral vasoconstriction lowers wrist temp 0.1–0.5°C",
        "3-axis acceleration: motion-artifact control + activity context",
    ], autofit=True)
    set_shape_text(s, "Text 4", "Prior ML on WESAD:")
    set_shape_text(s, "Text 5", [
        "Schmidt et al. SVM, handcrafted features: ~80% (3-class)",
        "TSFEL + Random Forest: ~85% (~4 MB feature vector)",
        "CNN on raw signals: ~87% (~2 MB, phone-tier)",
        "Unidirectional LSTM: ~88% (no future context)",
    ], autofit=True)

    # ---- Slide 6: Limitations ----
    print("Slide  6: Limitations of Existing Systems")
    s = slides[5]
    set_shape_text(s, "Title 1", "Limitations of Existing Systems")
    set_shape_text(s, "TextBox 9", [
        "No real-time wrist-only stress classifier deployed end-to-end on Wear OS",
        "Existing models too large (2–4 MB) for smartwatch memory & latency budgets",
        "Commercial stress scores: proprietary, unvalidated, not exposed to vehicles",
        "Vehicle driving modes (Comfort/Sport/Eco) are static, manually configured",
        "No published system unifies wearable + on-device ML + BLE + vehicle actuator",
    ], autofit=True)

    # ---- Slide 7: Proposed Solution ----
    print("Slide  7: Proposed Solution")
    s = slides[6]
    set_shape_text(s, "Title 1", "Proposed Solution")
    set_shape_text(s, "Text 2", "Five-Component Integrated System:")
    set_shape_text(s, "Text 3", [
        "Wearable Biosensing: Wear OS PPG/EDA/TEMP/ACC → 5-feature vector at 4 Hz",
        "AI Classifier: BiLSTM(64) + attention; 4 classes; 72 KB on-device TFLite",
        "BLE Communication: 12-byte GATT characteristic, 30 s cadence, mutual pairing",
        "Vehicle Integration: Arduino Nano 33 BLE central; 5-state FSM with WAITING failsafe",
        "Graduated Cabin Adaptation: RGB LED + PWM fan + audio per stress level",
    ], autofit=True)

    # ---- Slide 8: System Architecture (image swap) ----
    print("Slide  8: System Architecture (image swap)")
    s = slides[7]
    set_shape_text(s, "Title 1", "System Architecture")
    replace_picture(s, "Image 0", SYSTEM_OVERVIEW)
    set_shape_text(s, "TextBox 10", "End-to-End Pipeline: Wearable → BLE → Vehicle")

    # ---- Slide 9: Methodology & Data ----
    print("Slide  9: Methodology & Data")
    s = slides[8]
    set_shape_text(s, "Title 1", "Methodology & Data")
    set_shape_text(s, "Text 2", "WESAD Dataset:")
    set_shape_text(s, "Text 3", [
        "15 subjects (5 used); Empatica E4 wrist signals",
        "BVP / EDA / TEMP / ACC + Trier Social Stress Test protocol",
        "2,956 labelled 30×5 windows (7.5 s, 50% overlap); 80/20 stratified split",
    ], autofit=True)
    set_shape_text(s, "Text 4", "Pre-processing Pipeline:")
    set_shape_text(s, "Text 5", [
        "Resample all channels to 4 Hz; IBI extraction from 64 Hz BVP",
        "Per-subject EDA/IBI quantile thresholds → 4-class label mapping",
        "Z-score normalization fit on train only — no information leakage",
    ], autofit=True)

    # ---- Slide 10: Model Architecture (image swap) ----
    print("Slide 10: Model Architecture (image swap)")
    s = slides[9]
    set_shape_text(s, "Title 1", "Model Architecture")
    set_shape_text(s, "Text 2", "BiLSTM + Additive Attention:")
    set_shape_text(s, "Text 3", [
        "Input [1, 30, 5] — 30 timesteps × 5 features",
        "BiLSTM(64) → [1, 30, 128]  (35,840 params)",
        "Additive attention → [1, 128]  (129 params)",
        "Dense(64) → Dense(32) → Dense(4, softmax)  (10,468 params)",
        "Total: 160K params → 72 KB TFLite (4× quantization)",
    ], autofit=True)
    replace_picture(s, "Image 0", MODEL_ARCH)

    # ---- Slide 11: BLE Protocol & Cabin FSM ----
    print("Slide 11: BLE Protocol & Cabin FSM")
    s = slides[10]
    set_shape_text(s, "Title 1", "BLE Protocol & Cabin FSM")
    set_shape_text(s, "Text 2", "12-byte GATT Characteristic:")
    set_shape_text(s, "Text 3", [
        "Bytes 0–7: Timestamp (int64 LE, ms since epoch)",
        "Byte 8: Stress level (0=Calm, 1=Mild, 2=Moderate, 3=Critical)",
        "Byte 9: Confidence (uint8, 0–100%)",
        "Bytes 10–11: Reserved",
        "Cadence: 30 s emission; 60-s timeout → WAITING",
    ], autofit=True)
    set_shape_text(s, "Text 4", "5-State Cabin FSM:")
    set_shape_text(s, "Text 5", [
        "WAITING (failsafe) → yellow blink, fan off",
        "CALM → warm white, 30% fan, audio off",
        "MILD → neutral white, 50% fan",
        "MODERATE → blue, 75% fan, soft music",
        "CRITICAL → red blink, 100% fan, voice alert",
    ], autofit=True)

    # ---- Slide 12: Results & Evaluation ----
    print("Slide 12: Results & Evaluation")
    s = slides[11]
    set_shape_text(s, "Title 1", "Results & Evaluation")
    set_shape_text(s, "Text 2", "Test Performance (N=592 windows):")
    set_shape_text(s, "Text 3", [
        "Overall accuracy: 97.8% on held-out WESAD test split",
        "Per-class F1: Calm 0.979 | Mild 0.938 | Moderate 0.997",
        "Inference latency: 22.8 ms desktop; <50 ms target on Wear OS",
        "Model size: 72.2 KB (≤80 KB target ✓; 4× quantization)",
    ], autofit=True)
    set_shape_text(s, "Text 4", "Comparison with Prior Work:")
    set_shape_text(s, "Text 5", [
        "50–100× smaller than CNN/TSFEL baselines (2–4 MB)",
        "Higher accuracy on 4-class task vs. 3-class baselines (~85–88%)",
        "Only system validated end-to-end on deployment hardware",
        "All 5 integration scenarios pass (Calm/Mild/Moderate/Critical/Watch-removed)",
    ], autofit=True)

    # ---- Slide 13: Conclusion ----
    print("Slide 13: Conclusion")
    s = slides[12]
    set_shape_text(s, "Title 1", "Conclusion")
    set_shape_text(s, "Text 2", "Delivered Contributions:")
    set_shape_text(s, "Text 3", [
        "4-class wrist-only stress classifier: BiLSTM(64) + attention",
        "97.8% accuracy at 72 KB TFLite (exceeds ≥90% / ≤80 KB targets)",
        "Wear OS app with on-device inference <50 ms latency",
        "Documented 12-byte BLE GATT protocol — implementation-agnostic",
        "Arduino Nano 33 BLE cabin control with 5-state failsafe FSM",
    ], autofit=True)
    set_shape_text(s, "Text 4", "Future Work:")
    set_shape_text(s, "Text 5", [
        "Unification with AlcoWatch BAC system (ACN1408) on shared hardware",
        "Real-world driving simulator / test-track validation",
        "Per-driver personalization via 5-min baseline calibration",
        "Knowledge distillation to <25 KB student model",
    ], autofit=True)
    # Practical-implementation proof callout — fits in the gap between Text 5
    # (ends at top≈6.16in) and the slide-number footer (top=6.95in).
    add_centered_textbox(
        s, Inches(0.5), Inches(6.20), Inches(12.3), Inches(0.50),
        f"Open-Source Implementation:  {GITHUB_URL}",
        font_size=14, bold=True, color=(30, 30, 90),
    )

    # ---- Slide 14: References ----
    print("Slide 14: References")
    s = slides[13]
    set_shape_text(s, "Title 1", "References")
    set_shape_text(s, "Content Placeholder 4", [
        "Schmidt P. et al. (2018). WESAD: Multimodal Dataset for Wearable Stress Detection. ACM ICMI.",
        "Graves A. & Schmidhuber J. (2005). Bidirectional LSTM Phoneme Classification. Neural Networks 18.",
        "Vaswani A. et al. (2017). Attention Is All You Need. NeurIPS 30.",
        "Sharma N. & Gedeon T. (2012). Stress Recognition: A Survey. CMPB 108, 1287–1301.",
        "Can Y.S., Arnrich B., Ersoy C. (2019). Stress Detection with Wearables. JBI 92.",
        "Barandas M. et al. (2020). TSFEL: Time Series Feature Library. SoftwareX 11.",
        "Koldijk S. et al. (2014). The SWELL Knowledge Work Dataset. ACM ICMI.",
        "Google (2023). TensorFlow Lite: On-Device ML. tensorflow.org/lite",
    ], autofit=True)

    print("Slide 15: Thank You (+ GitHub link callout)")
    s = slides[14]
    # Rectangle 2 ("Thank You") sits at top=3.25, h=2.42 → bottom=5.67.
    # Slide-number footer at top=6.95. We have 1.28in of space underneath.
    add_centered_textbox(
        s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.55),
        "Code, Model & Firmware — Open Source",
        font_size=18, bold=False, color=(80, 80, 80),
    )
    add_centered_textbox(
        s, Inches(0.5), Inches(6.30), Inches(12.3), Inches(0.55),
        GITHUB_URL,
        font_size=22, bold=True, color=(30, 30, 90),
    )

    p.save(str(OUTPUT))
    print()
    print(f"OK Saved: {OUTPUT}")
    print(f"   Size: {OUTPUT.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
